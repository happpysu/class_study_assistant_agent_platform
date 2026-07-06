"""上传资料的文本抽取与切片。"""
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

TEXT_SUFFIXES = {".txt", ".md", ".markdown", ".py", ".c", ".cpp", ".java", ".json", ".csv"}

CHUNK_SIZE = 600
CHUNK_OVERLAP = 100


def extract_text(path: Path) -> str:
    """按文件类型抽取纯文本；不支持的类型返回空串（仅存档不检索）。"""
    suffix = path.suffix.lower()
    try:
        if suffix in TEXT_SUFFIXES:
            return path.read_text(encoding="utf-8", errors="ignore")
        if suffix == ".pdf":
            return _extract_pdf(path)
        if suffix == ".docx":
            return _extract_docx(path)
        if suffix == ".pptx":
            return _extract_pptx(path)
    except Exception as exc:  # 抽取失败不应阻断上传
        logger.warning("文本抽取失败 %s: %s", path.name, exc)
    return ""


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)


def _extract_docx(path: Path) -> str:
    """Word 文档：段落 + 表格单元格文本。"""
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    """PPT 幻灯片：逐页提取所有文本框内容，标注页码便于引用定位。"""
    from pptx import Presentation

    pres = Presentation(str(path))
    pages = []
    for idx, slide in enumerate(pres.slides, start=1):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            pages.append(f"[第{idx}页]\n" + "\n".join(texts))
    return "\n\n".join(pages)


def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """滑动窗口切片，保留少量重叠避免语义断裂。"""
    cleaned = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not cleaned:
        return []
    if size <= overlap:
        raise ValueError("size 必须大于 overlap")
    chunks = []
    start = 0
    while start < len(cleaned):
        chunks.append(cleaned[start : start + size])
        start += size - overlap
    return chunks
