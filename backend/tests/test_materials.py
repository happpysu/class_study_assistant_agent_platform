from .conftest import create_course

SAMPLE_TEXT = (
    "第一章 极限与连续。极限是微积分的基础概念，函数在某点的极限描述其趋近行为。\n"
    "第二章 导数与微分。导数刻画函数变化率，可导必连续。\n"
    "第三章 积分。定积分表示曲边梯形面积，牛顿-莱布尼茨公式联系微分与积分。"
)


def _upload(client, headers, course_id, filename="note.txt", mtype="notes"):
    return client.post(
        f"/api/courses/{course_id}/materials",
        files={"file": (filename, SAMPLE_TEXT.encode(), "text/plain")},
        data={"mtype": mtype, "description": "复习笔记"},
        headers=headers,
    )


def test_upload_list_search_delete(client, auth_headers):
    course_id = create_course(client, auth_headers)

    resp = _upload(client, auth_headers, course_id)
    assert resp.status_code == 201, resp.text
    material_id = resp.json()["id"]

    resp = client.get(f"/api/courses/{course_id}/materials", headers=auth_headers)
    assert len(resp.json()) == 1

    resp = client.get(
        f"/api/courses/{course_id}/materials",
        params={"mtype": "courseware"},
        headers=auth_headers,
    )
    assert resp.json() == []

    resp = client.get(
        f"/api/courses/{course_id}/materials/search",
        params={"q": "导数"},
        headers=auth_headers,
    )
    assert resp.status_code == 200
    hits = resp.json()
    assert hits and "导数" in hits[0]["excerpt"]

    resp = client.get(
        f"/api/materials/{material_id}/download", headers=auth_headers
    )
    assert resp.status_code == 200

    assert (
        client.delete(
            f"/api/materials/{material_id}", headers=auth_headers
        ).status_code
        == 204
    )
    resp = client.get(f"/api/courses/{course_id}/materials", headers=auth_headers)
    assert resp.json() == []


def test_upload_invalid_type(client, auth_headers):
    course_id = create_course(client, auth_headers)
    resp = _upload(client, auth_headers, course_id, mtype="bad-type")
    assert resp.status_code == 422


def test_upload_docx_parsed_and_searchable(client, auth_headers):
    import io

    from docx import Document

    course_id = create_course(client, auth_headers)
    doc = Document()
    doc.add_paragraph("第四章 多元函数微分学。偏导数是多元函数对单个自变量的导数。")
    buf = io.BytesIO()
    doc.save(buf)

    resp = client.post(
        f"/api/courses/{course_id}/materials",
        files={
            "file": (
                "chap4.docx",
                buf.getvalue(),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
        data={"mtype": "courseware"},
        headers=auth_headers,
    )
    assert resp.status_code == 201, resp.text

    hits = client.get(
        f"/api/courses/{course_id}/materials/search",
        params={"q": "偏导数"},
        headers=auth_headers,
    ).json()
    assert hits and hits[0]["material_name"] == "chap4.docx"


def test_upload_pptx_parsed_and_searchable(client, auth_headers):
    import io

    from pptx import Presentation

    course_id = create_course(client, auth_headers)
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.shapes.title.text = "第五章 重积分"
    slide.placeholders[1].text = "二重积分的几何意义是曲顶柱体的体积。"
    buf = io.BytesIO()
    pres.save(buf)

    resp = client.post(
        f"/api/courses/{course_id}/materials",
        files={
            "file": (
                "chap5.pptx",
                buf.getvalue(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
        data={"mtype": "courseware"},
        headers=auth_headers,
    )
    assert resp.status_code == 201, resp.text

    hits = client.get(
        f"/api/courses/{course_id}/materials/search",
        params={"q": "二重积分"},
        headers=auth_headers,
    ).json()
    assert hits and hits[0]["material_name"] == "chap5.pptx"
    assert "第1页" in hits[0]["excerpt"] or "二重积分" in hits[0]["excerpt"]
